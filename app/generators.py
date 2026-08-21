from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from . import design
from .design import Palette, palette_from_profile
from .models import DocumentProfile

FOOTER_NOTE = "Multi-Agent Document & PPT POC"


def _font_name(profile: dict[str, Any], fallback: str) -> str:
    names = profile.get("style", {}).get("font_names", [])
    return names[0] if names else fallback


def _font_size(profile: dict[str, Any], fallback: float) -> float:
    sizes = profile.get("style", {}).get("font_sizes", [])
    return float(sizes[0]) if sizes else fallback


def _clamp(value: float, low: float, high: float) -> float:
    return max(low, min(high, value))


def _today() -> str:
    return datetime.now(timezone.utc).strftime("%d %B %Y")


# ---------------------------------------------------------------------------
# DOCX generation
# ---------------------------------------------------------------------------


def _style_run(
    run: Any,
    *,
    font: str,
    size: float,
    color: str,
    bold: bool = False,
    italic: bool = False,
    caps: bool = False,
) -> Any:
    from docx.shared import Pt, RGBColor

    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(design.clean_hex(color))
    run.font.bold = bold
    run.font.italic = italic
    run.font.all_caps = caps
    return run


def _spacing(paragraph: Any, *, before: float = 0, after: float = 6, line: float | None = None) -> Any:
    from docx.shared import Pt

    fmt = paragraph.paragraph_format
    fmt.space_before = Pt(before)
    fmt.space_after = Pt(after)
    if line is not None:
        fmt.line_spacing = line
    return paragraph


def _new_paragraph(document: Any, style_name: str | None = None) -> Any:
    """Add a paragraph, degrading gracefully when a template lacks the style."""
    if style_name and style_name in {style.name for style in document.styles}:
        return document.add_paragraph(style=style_name)
    return document.add_paragraph()


def _panel_table(document: Any, fill: str, *, padding: int = 260) -> Any:
    """A single borderless, shaded cell used for cover panels and callouts."""
    table = document.add_table(rows=1, cols=1)
    design.clear_table_borders(table)
    design.set_table_width_pct(table, 100)
    cell = table.cell(0, 0)
    design.shade_cell(cell, fill)
    design.cell_margins(cell, top=padding, start=padding, bottom=padding, end=padding)
    return cell


def _cell_paragraph(cell: Any, first: bool = False) -> Any:
    return cell.paragraphs[0] if first else cell.add_paragraph()


def _apply_base_styles(document: Any, palette: Palette, body_font: str, base_size: float) -> None:
    """Restyle the built-in styles when we own the document (no uploaded template)."""
    from docx.oxml.ns import qn
    from docx.shared import Pt, RGBColor

    normal = document.styles["Normal"]
    normal.font.name = body_font
    normal.font.size = Pt(base_size)
    normal.font.color.rgb = RGBColor.from_string(palette.body)
    rfonts = normal.element.get_or_add_rPr().get_or_add_rFonts()
    for attribute in ("w:ascii", "w:hAnsi", "w:cs", "w:eastAsia"):
        rfonts.set(qn(attribute), body_font)
    normal.paragraph_format.space_after = Pt(7)
    normal.paragraph_format.line_spacing = 1.22

    available = {style.name for style in document.styles}
    heading_specs = {
        "Heading 1": (base_size + 5.5, palette.ink),
        "Heading 2": (base_size + 2.5, palette.accent),
        "Heading 3": (base_size + 0.5, palette.accent_deep),
    }
    for name, (size, color) in heading_specs.items():
        if name not in available:
            continue
        style = document.styles[name]
        style.font.name = body_font
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(color)
        style.paragraph_format.space_before = Pt(14)
        style.paragraph_format.space_after = Pt(5)
        style.paragraph_format.keep_with_next = True

    if "List Bullet" in available:
        bullet = document.styles["List Bullet"]
        bullet.font.name = body_font
        bullet.font.size = Pt(base_size)
        bullet.font.color.rgb = RGBColor.from_string(palette.body)
        bullet.paragraph_format.space_after = Pt(4)
        bullet.paragraph_format.line_spacing = 1.18


def _docx_cover(document: Any, content: dict[str, Any], palette: Palette, font: str, base_size: float) -> None:
    cell = _panel_table(document, palette.accent_deep, padding=300)

    eyebrow = _cell_paragraph(cell, first=True)
    _spacing(eyebrow, after=8)
    eyebrow_run = _style_run(
        eyebrow.add_run("Multi-agent research briefing"),
        font=font,
        size=max(7.5, base_size - 2.5),
        color=palette.on_accent_muted,
        bold=True,
        caps=True,
    )
    design.letter_spacing(eyebrow_run, 46)

    heading = _cell_paragraph(cell)
    _spacing(heading, after=6, line=1.05)
    _style_run(
        heading.add_run(content.get("title", "AI Strategy Proposal")),
        font=font,
        size=base_size + 15,
        color=palette.on_accent,
        bold=True,
    )

    keyline = _cell_paragraph(cell)
    _spacing(keyline, before=2, after=8)
    design.paragraph_border(keyline, "bottom", palette.on_accent_muted, size=12, space=2)

    subtitle = content.get("subtitle")
    if subtitle:
        line = _cell_paragraph(cell)
        _spacing(line, after=0, line=1.18)
        _style_run(line.add_run(str(subtitle)), font=font, size=base_size + 0.5, color=palette.on_accent_muted)

    _spacing(document.add_paragraph(), after=2)

    facts = [
        ("Prepared by", "Supervisor agent workflow"),
        ("Generated", _today()),
        ("Evidence", f"{len(content.get('sources') or [])} cited source(s)"),
    ]
    strip = document.add_table(rows=2, cols=len(facts))
    design.clear_table_borders(strip)
    design.set_table_width_pct(strip, 100)
    for index, (label, value) in enumerate(facts):
        label_cell = strip.cell(0, index)
        design.set_cell_width_pct(label_cell, 100 / len(facts))
        design.cell_margins(label_cell, top=40, start=0, bottom=0, end=120)
        label_paragraph = label_cell.paragraphs[0]
        _spacing(label_paragraph, after=1)
        label_run = _style_run(
            label_paragraph.add_run(label),
            font=font,
            size=max(7.0, base_size - 3.0),
            color=palette.muted,
            bold=True,
            caps=True,
        )
        design.letter_spacing(label_run, 34)

        value_cell = strip.cell(1, index)
        design.set_cell_width_pct(value_cell, 100 / len(facts))
        design.cell_margins(value_cell, top=0, start=0, bottom=40, end=120)
        value_paragraph = value_cell.paragraphs[0]
        _spacing(value_paragraph, after=0)
        _style_run(value_paragraph.add_run(value), font=font, size=base_size - 0.5, color=palette.ink, bold=True)
        design.paragraph_border(value_paragraph, "bottom", palette.rule, size=6, space=4)

    _spacing(document.add_paragraph(), after=10)


def _docx_section_heading(
    document: Any,
    number: int | None,
    text: str,
    palette: Palette,
    font: str,
    base_size: float,
) -> Any:
    paragraph = _new_paragraph(document, "Heading 1")
    _spacing(paragraph, before=16, after=6)
    design.keep_with_next(paragraph)
    size = base_size + 5.5
    if number is not None:
        badge = _style_run(paragraph.add_run(f"{number:02d}"), font=font, size=size, color=palette.accent_soft, bold=True)
        design.letter_spacing(badge, 20)
        _style_run(paragraph.add_run("   "), font=font, size=size, color=palette.accent_soft)
    _style_run(paragraph.add_run(text), font=font, size=size, color=palette.ink, bold=True)
    design.paragraph_border(paragraph, "bottom", palette.accent, size=10, space=6)
    return paragraph


def _docx_body(document: Any, text: str, palette: Palette, font: str, base_size: float) -> Any:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    _spacing(paragraph, before=2, after=8, line=1.28)
    _style_run(paragraph.add_run(text), font=font, size=base_size, color=palette.body)
    return paragraph


def _docx_bullets(document: Any, bullets: list[str], palette: Palette, font: str, base_size: float) -> None:
    from docx.shared import Inches

    for bullet in bullets:
        paragraph = _new_paragraph(document, "List Bullet")
        _spacing(paragraph, before=0, after=4, line=1.2)
        paragraph.paragraph_format.left_indent = Inches(0.32)
        _style_run(paragraph.add_run(str(bullet)), font=font, size=base_size, color=palette.body)


def _docx_callout(document: Any, label: str, text: str, palette: Palette, font: str, base_size: float) -> None:
    cell = _panel_table(document, palette.tint, padding=200)
    design.cell_borders(
        cell,
        {
            "left": {"sz": 30, "color": palette.accent},
            "top": {"sz": 4, "color": palette.tint_strong},
            "bottom": {"sz": 4, "color": palette.tint_strong},
            "right": {"sz": 4, "color": palette.tint_strong},
        },
    )
    heading = _cell_paragraph(cell, first=True)
    _spacing(heading, after=4)
    label_run = _style_run(
        heading.add_run(label),
        font=font,
        size=max(7.0, base_size - 3.0),
        color=palette.accent,
        bold=True,
        caps=True,
    )
    design.letter_spacing(label_run, 40)

    body = _cell_paragraph(cell)
    _spacing(body, after=0, line=1.26)
    _style_run(body.add_run(text), font=font, size=base_size + 0.5, color=palette.ink)

    _spacing(document.add_paragraph(), after=6)


def _docx_table(
    document: Any,
    headers: list[str],
    rows: list[list[str]],
    widths: list[float],
    palette: Palette,
    font: str,
    base_size: float,
) -> None:
    table = document.add_table(rows=1, cols=len(headers))
    design.clear_table_borders(table)
    design.set_table_width_pct(table, 100)

    design.repeat_header_row(table.rows[0])
    design.keep_row_together(table.rows[0])
    for index, title in enumerate(table.rows[0].cells):
        design.shade_cell(title, palette.accent)
        design.set_cell_width_pct(title, widths[index])
        design.cell_margins(title, top=90, start=130, bottom=90, end=130)
        paragraph = title.paragraphs[0]
        _spacing(paragraph, after=0)
        header_run = _style_run(
            paragraph.add_run(headers[index]),
            font=font,
            size=max(7.0, base_size - 2.5),
            color=palette.on_accent,
            bold=True,
            caps=True,
        )
        design.letter_spacing(header_run, 34)

    for row_index, values in enumerate(rows):
        row = table.add_row()
        design.keep_row_together(row)
        cells = row.cells
        fill = palette.surface if row_index % 2 else design.WHITE
        for column, value in enumerate(values[: len(headers)]):
            cell = cells[column]
            design.shade_cell(cell, fill)
            design.set_cell_width_pct(cell, widths[column])
            design.cell_margins(cell, top=80, start=130, bottom=80, end=130)
            design.cell_borders(cell, {"bottom": {"sz": 4, "color": palette.rule}})
            paragraph = cell.paragraphs[0]
            _spacing(paragraph, after=0, line=1.16)
            _style_run(
                paragraph.add_run(str(value)),
                font=font,
                size=base_size - 0.5,
                color=palette.muted if column == 0 else palette.body,
                bold=column == 0,
            )

    _spacing(document.add_paragraph(), after=6)


def _docx_footer(document: Any, content: dict[str, Any], palette: Palette, font: str, base_size: float) -> None:
    from docx.enum.text import WD_TAB_ALIGNMENT
    from docx.shared import Inches

    size = max(7.0, base_size - 3.0)
    for section in document.sections:
        paragraph = section.footer.paragraphs[0] if section.footer.paragraphs else section.footer.add_paragraph()
        paragraph.text = ""
        _spacing(paragraph, before=4, after=0)
        design.paragraph_border(paragraph, "top", palette.rule, size=6, space=6)
        usable = section.page_width.inches - section.left_margin.inches - section.right_margin.inches
        paragraph.paragraph_format.tab_stops.add_tab_stop(Inches(usable), WD_TAB_ALIGNMENT.RIGHT)
        label = f"{content.get('title', 'Generated document')}  {design.DOT_CHAR}  {FOOTER_NOTE}"
        _style_run(paragraph.add_run(label), font=font, size=size, color=palette.muted)
        _style_run(paragraph.add_run("\t"), font=font, size=size, color=palette.muted)
        _style_run(paragraph.add_run("Page "), font=font, size=size, color=palette.muted)
        design.add_field(paragraph, "PAGE")


def generate_docx(content: dict[str, Any], profile: dict[str, Any], output_path: Path, template_path: Path | None = None) -> Path:
    from docx import Document
    from docx.shared import Inches

    template_in_use = bool(template_path and template_path.exists())
    document = Document(str(template_path)) if template_in_use else Document()

    palette = palette_from_profile(profile)
    body_font = _font_name(profile, "Calibri")
    base_size = _clamp(_font_size(profile, 10.5), 9.5, 12.0)

    if not template_in_use:
        margins = profile.get("style", {}).get("margins", {})
        for section in document.sections:
            section.top_margin = Inches(float(margins.get("top", 0.9)))
            section.bottom_margin = Inches(float(margins.get("bottom", 0.9)))
            section.left_margin = Inches(float(margins.get("left", 0.95)))
            section.right_margin = Inches(float(margins.get("right", 0.95)))
        _apply_base_styles(document, palette, body_font, base_size)
        # Headers and footers of an uploaded template are part of its fidelity
        # contract, so page furniture is only added to documents we own.
        _docx_footer(document, content, palette, body_font, base_size)
    elif document.paragraphs:
        # Keep the supplied body, tables, images, sections, headers, footers, and
        # custom styles intact; the styled deliverable starts on a fresh page.
        document.add_page_break()

    _docx_cover(document, content, palette, body_font, base_size)

    _docx_section_heading(document, None, "Executive Summary", palette, body_font, base_size)
    summary = content.get("executive_summary", "")
    if summary:
        _docx_callout(document, "At a glance", summary, palette, body_font, base_size)

    for index, section in enumerate(content.get("sections", []), start=1):
        _docx_section_heading(document, index, section.get("heading", "Section"), palette, body_font, base_size)
        if section.get("body"):
            _docx_body(document, section["body"], palette, body_font, base_size)
        bullets = [str(item) for item in section.get("bullets", []) if str(item).strip()]
        if bullets:
            _docx_bullets(document, bullets, palette, body_font, base_size)

    _docx_section_heading(document, None, "Sources and Traceability", palette, body_font, base_size)
    sources = content.get("sources") or []
    if sources:
        _docx_body(
            document,
            "Every claim in this document is traceable to the evidence the Research Agent returned for this run.",
            palette,
            body_font,
            base_size,
        )
        _docx_table(
            document,
            ["#", "Source", "Reference"],
            [
                [
                    f"{index:02d}",
                    str(source.get("title") or "Untitled source"),
                    str(source.get("url") or source.get("snippet") or "No reference recorded"),
                ]
                for index, source in enumerate(sources, start=1)
            ],
            [8.0, 42.0, 50.0],
            palette,
            body_font,
            base_size,
        )
    else:
        _docx_callout(
            document,
            "Evidence status",
            "No external sources were returned for this run; the content is labelled as deterministic fallback material.",
            palette,
            body_font,
            base_size,
        )

    claim_citations = content.get("claim_citations") or []
    if claim_citations:
        _docx_section_heading(document, None, "Claim-Level Citation Map", palette, body_font, base_size)
        _docx_table(
            document,
            ["#", "Claim", "Cited sources"],
            [
                [
                    f"{index:02d}",
                    str(citation.get("claim", "")),
                    "\n".join(str(url) for url in citation.get("source_urls", [])) or "Not cited",
                ]
                for index, citation in enumerate(claim_citations, start=1)
            ],
            [8.0, 52.0, 40.0],
            palette,
            body_font,
            base_size,
        )

    document.core_properties.title = content.get("title", "AI Strategy Proposal")
    document.core_properties.subject = "Generated by the Multi-Agent Document and PPT POC"
    document.core_properties.author = "Supervisor agent workflow"
    document.core_properties.category = "Research-backed proposal"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    document.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# PPTX generation
# ---------------------------------------------------------------------------


class SlideCanvas:
    """Resolution-independent slide geometry derived from the deck's slide size."""

    def __init__(self, width_inches: float, height_inches: float) -> None:
        self.width = width_inches
        self.height = height_inches
        self.margin = _clamp(width_inches * 0.062, 0.45, 0.9)
        self.rail = _clamp(width_inches * 0.013, 0.09, 0.18)
        self.content_width = width_inches - (self.margin * 2)
        self.title_top = _clamp(height_inches * 0.115, 0.55, 1.05)
        self.rule_top = self.title_top + _clamp(height_inches * 0.128, 0.85, 1.25)
        self.body_top = self.rule_top + 0.28
        self.footer_top = height_inches - _clamp(height_inches * 0.085, 0.5, 0.8)
        self.body_height = max(1.0, self.footer_top - self.body_top - 0.18)
        self.accent_rule_width = _clamp(width_inches * 0.16, 1.1, 2.0)

    @property
    def title_scale(self) -> float:
        return _clamp(self.width / 10.0, 0.85, 1.25)


def _rgb(color: str) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(design.clean_hex(color))


def _rect(slide: Any, left: float, top: float, width: float, height: float, fill: str, *, behind: bool = True) -> Any:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    design.flat_shape(shape)
    if behind:
        design.send_to_back(shape)
    return shape


def _textbox(slide: Any, left: float, top: float, width: float, height: float) -> Any:
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return box


def _style_pptx_run(run: Any, *, font: str, size: float, color: str, bold: bool = False) -> Any:
    from pptx.util import Pt

    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return run


def _single_line(box: Any, text: str, *, font: str, size: float, color: str, bold: bool = False, align: Any = None) -> Any:
    paragraph = box.text_frame.paragraphs[0]
    design.clear_bullet(paragraph)
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    _style_pptx_run(run, font=font, size=size, color=color, bold=bold)
    return paragraph


BODY_SIZE_STEPS = (24.0, 22.0, 20.0, 19.0, 18.0, 17.0, 16.0, 15.0, 14.0, 13.0, 12.0)
DETAIL_SIZE_RATIO = 0.76
_HANGING_INDENT_INCHES = 0.23

Bullet = tuple[str, str]


def _normalize_bullets(bullets: list[Any]) -> list[Bullet]:
    """Accept plain strings or ``{"text": ..., "detail": ...}`` two-part bullets."""
    normalized: list[Bullet] = []
    for item in bullets:
        if isinstance(item, dict):
            text = str(item.get("text", "")).strip()
            detail = str(item.get("detail", "") or "").strip()
        else:
            text, detail = str(item).strip(), ""
        if text or detail:
            normalized.append((text or detail, "" if text == detail else detail))
    return normalized


def _wrapped_lines(text: str, size: float, width_inches: float) -> int:
    """Estimate how many rendered lines a bullet needs at a given point size."""
    if not text:
        return 0
    average_glyph = size * 0.5 / 72  # inches; a workable mean for proportional faces
    per_line = max(8, int((width_inches - _HANGING_INDENT_INCHES) / average_glyph))
    return max(1, -(-len(text) // per_line))


def _fit_body(bullets: list[Bullet], width_inches: float, height_inches: float, scale: float) -> tuple[float, float]:
    """Pick the largest readable size that fits, plus the gap that distributes the block."""
    budget = height_inches * 0.92
    ceiling = BODY_SIZE_STEPS[0] * scale
    for step in BODY_SIZE_STEPS:
        size = round(min(step, ceiling), 1)
        detail_size = size * DETAIL_SIZE_RATIO
        text_height = sum(
            _wrapped_lines(text, size, width_inches) * (size * 1.2 / 72)
            + _wrapped_lines(detail, detail_size, width_inches) * (detail_size * 1.2 / 72)
            for text, detail in bullets
        )
        minimum_gaps = max(0, len(bullets) - 1) * (9 / 72)
        if text_height + minimum_gaps <= budget or step == BODY_SIZE_STEPS[-1]:
            leftover = max(0.0, budget - text_height - minimum_gaps)
            slots = max(1, len(bullets) - 1)
            gap = 9.0 + leftover * 72 / slots if len(bullets) > 1 else 0.0
            return size, round(min(gap, 30.0), 1)
    return BODY_SIZE_STEPS[-1], 9.0


def _title_font_size(title: str, canvas: SlideCanvas) -> float:
    length = len(title)
    size = 30.0 if length <= 30 else 26.0 if length <= 48 else 22.0
    return round(size * canvas.title_scale, 1)


def _apply_chrome(
    slide: Any,
    canvas: SlideCanvas,
    palette: Palette,
    font: str,
    index: int,
    total: int,
    deck_title: str,
    dark: bool = False,
) -> None:
    from pptx.enum.text import PP_ALIGN

    _rect(slide, 0, 0, canvas.rail, canvas.height, palette.on_accent_muted if dark else palette.accent, behind=False)

    rule_color = design.mix(palette.on_accent, palette.accent_deep, 0.55) if dark else palette.rule
    _rect(slide, canvas.margin, canvas.footer_top - 0.06, canvas.content_width, 0.012, rule_color, behind=False)

    footer = _textbox(slide, canvas.margin, canvas.footer_top, canvas.content_width * 0.72, 0.28)
    _single_line(
        footer,
        f"{deck_title}  {design.DOT_CHAR}  {FOOTER_NOTE}",
        font=font,
        size=9.0,
        color=palette.on_accent_muted if dark else palette.muted,
    )

    number = _textbox(slide, canvas.width - canvas.margin - 1.1, canvas.footer_top, 1.1, 0.28)
    _single_line(
        number,
        f"{index:02d} / {total:02d}",
        font=font,
        size=9.5,
        color=palette.on_accent if dark else palette.accent,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def _fill_bullets(
    shape: Any,
    bullets: list[Any],
    palette: Palette,
    canvas: SlideCanvas,
    font: str,
    dark: bool = False,
) -> None:
    """Lay bullets into a shape, sized and spaced to occupy the space available."""
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Pt

    items = _normalize_bullets(bullets) or [
        ("Content will be refined through the conversational editing workflow.", "")
    ]
    width, height = canvas.content_width, canvas.body_height
    try:
        width = shape.width / 914400
        height = shape.height / 914400
    except (TypeError, AttributeError):  # pragma: no cover - template dependent
        pass
    size, gap = _fit_body(items, width, height, _clamp(canvas.height / 7.5, 0.85, 1.15))

    text_color = palette.on_accent if dark else palette.body
    detail_color = palette.on_accent_muted if dark else palette.muted
    bullet_color = palette.on_accent_muted if dark else palette.accent
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for position, (text, detail) in enumerate(items):
        paragraph = frame.paragraphs[0] if position == 0 else frame.add_paragraph()
        design.set_bullet(paragraph, bullet_color)
        paragraph.space_before = Pt(0 if position == 0 else gap)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.2
        run = paragraph.add_run()
        run.text = text
        _style_pptx_run(run, font=font, size=size, color=text_color)
        if detail:
            paragraph.add_line_break()
            detail_run = paragraph.add_run()
            detail_run.text = detail
            _style_pptx_run(detail_run, font=font, size=round(size * DETAIL_SIZE_RATIO, 1), color=detail_color)


def _eyebrow(
    slide: Any,
    canvas: SlideCanvas,
    palette: Palette,
    font: str,
    text: str,
    title_top: float,
    left: float,
    dark: bool = False,
) -> None:
    """Small tracked-out section label seated just above the slide title."""
    top = max(0.06, title_top - 0.30)
    if top + 0.24 > title_top + 0.02:
        return
    box = _textbox(slide, left, top, canvas.content_width, 0.24)
    _single_line(
        box,
        text.upper(),
        font=font,
        size=9.5,
        color=palette.on_accent_muted if dark else palette.muted,
        bold=True,
    )


def _blank_layout(presentation: Any) -> Any:
    layouts = list(presentation.slide_layouts)
    if not layouts:
        raise ValueError("The presentation template exposes no slide layouts")
    for layout in layouts:
        if (layout.name or "").strip().lower() == "blank":
            return layout
    return min(layouts, key=lambda item: len(item.placeholders))


def _compose_slide(
    presentation: Any,
    canvas: SlideCanvas,
    palette: Palette,
    font: str,
    title_text: str,
    bullets: list[str],
    index: int,
    total: int,
    deck_title: str,
    eyebrow: str,
    dark: bool = False,
) -> Any:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    slide = presentation.slides.add_slide(_blank_layout(presentation))
    if dark:
        _rect(slide, 0, 0, canvas.width, canvas.height, palette.accent_deep)

    _eyebrow(slide, canvas, palette, font, eyebrow, canvas.title_top, canvas.margin, dark)

    title_box = _textbox(
        slide,
        canvas.margin,
        canvas.title_top,
        canvas.content_width,
        max(0.5, canvas.rule_top - canvas.title_top - 0.12),
    )
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    title_paragraph = _single_line(
        title_box,
        title_text,
        font=font,
        size=_title_font_size(title_text, canvas),
        color=palette.on_accent if dark else palette.ink,
        bold=True,
        align=PP_ALIGN.LEFT,
    )
    title_paragraph.line_spacing = 1.04

    _rect(
        slide,
        canvas.margin,
        canvas.rule_top,
        canvas.accent_rule_width,
        0.055,
        palette.on_accent_muted if dark else palette.accent,
        behind=False,
    )

    body = _textbox(slide, canvas.margin, canvas.body_top, canvas.content_width, canvas.body_height)
    _fill_bullets(body, bullets, palette, canvas, font, dark)

    _apply_chrome(slide, canvas, palette, font, index, total, deck_title, dark)
    return slide


def _restyle_template_slide(
    slide: Any,
    canvas: SlideCanvas,
    palette: Palette,
    font: str,
    title_text: str,
    bullets: list[str],
    index: int,
    total: int,
    deck_title: str,
    eyebrow: str,
) -> Any:
    """Fill an existing template slide without disturbing its placeholder geometry."""
    from pptx.enum.text import PP_ALIGN

    title = slide.shapes.title
    if title is None:
        title = next((shape for shape in slide.placeholders if getattr(shape, "has_text_frame", False)), None)
    if title is None:
        title = _textbox(slide, canvas.margin, canvas.title_top, canvas.content_width, 0.9)
    title.text_frame.clear()
    _single_line(
        title,
        title_text,
        font=font,
        size=_title_font_size(title_text, canvas),
        color=palette.ink,
        bold=True,
        align=PP_ALIGN.LEFT,
    )

    title_idx = getattr(title.placeholder_format, "idx", -1) if getattr(title, "is_placeholder", False) else -1
    body = next(
        (
            shape
            for shape in slide.placeholders
            if getattr(shape, "has_text_frame", False) and getattr(shape.placeholder_format, "idx", -1) != title_idx
        ),
        None,
    )
    if body is None:
        body = next((shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape is not title), None)
    if body is None:
        body = _textbox(slide, canvas.margin, canvas.body_top, canvas.content_width, canvas.body_height)
    _fill_bullets(body, bullets, palette, canvas, font)

    # A section label and a short accent rule tie the reused template slide into
    # the same visual system as the generated slides, without moving placeholders.
    rule_top, rule_left, title_top = canvas.rule_top, canvas.margin, canvas.title_top
    try:
        title_top = title.top / 914400
        rule_top = (title.top + title.height) / 914400 + 0.06
        rule_left = max(canvas.rail + 0.08, title.left / 914400)
    except (TypeError, AttributeError):  # pragma: no cover - template dependent
        pass
    _eyebrow(slide, canvas, palette, font, eyebrow, title_top, rule_left)
    if rule_top < canvas.footer_top - 0.3:
        _rect(slide, rule_left, rule_top, canvas.accent_rule_width, 0.055, palette.accent, behind=False)

    _apply_chrome(slide, canvas, palette, font, index, total, deck_title)
    return slide


def generate_pptx(content: dict[str, Any], profile: dict[str, Any], output_path: Path, template_path: Path | None = None) -> Path:
    from pptx import Presentation

    presentation = Presentation(str(template_path)) if template_path and template_path.exists() else Presentation()
    palette = palette_from_profile(profile)
    font = _font_name(profile, "Calibri")
    canvas = SlideCanvas(presentation.slide_width / 914400, presentation.slide_height / 914400)
    deck_title = content.get("title", "AI Strategy Presentation")

    body_slides = content.get("slides", [])[:11]
    total = len(body_slides) + 1
    template_slide_count = len(presentation.slides)

    for index, slide_content in enumerate(body_slides):
        title_text = slide_content.get("title", "Untitled")
        bullets = [str(item) for item in slide_content.get("bullets", []) if str(item).strip()]
        if index < template_slide_count:
            _restyle_template_slide(
                presentation.slides[index],
                canvas,
                palette,
                font,
                title_text,
                bullets,
                index + 1,
                total,
                deck_title,
                eyebrow=f"Section {index + 1:02d}",
            )
        else:
            _compose_slide(
                presentation,
                canvas,
                palette,
                font,
                title_text,
                bullets,
                index + 1,
                total,
                deck_title,
                eyebrow=f"Section {index + 1:02d}",
            )

    sources = content.get("sources") or []
    source_bullets: list[Any] = [
        {
            "text": str(source.get("title") or "Untitled source"),
            "detail": str(source.get("url") or source.get("snippet") or "No reference recorded"),
        }
        for source in sources
    ] or ["No external sources were returned; the demo content is labelled accordingly."]
    source_index = len(body_slides)
    if source_index < template_slide_count:
        _restyle_template_slide(
            presentation.slides[source_index],
            canvas,
            palette,
            font,
            "Sources & Traceability",
            source_bullets,
            total,
            total,
            deck_title,
            eyebrow="Evidence",
        )
    else:
        _compose_slide(
            presentation,
            canvas,
            palette,
            font,
            "Sources & Traceability",
            source_bullets,
            total,
            total,
            deck_title,
            eyebrow="Evidence",
            dark=True,
        )

    source_notes = "[Sources]\n" + "\n".join(f"- {source.get('title')}: {source.get('url')}" for source in sources)
    claim_citations = content.get("claim_citations", [])
    if claim_citations:
        source_notes += "\n[Claim citations]\n" + "\n".join(
            f"- {item.get('claim', '')}: {', '.join(item.get('source_urls', []))}" for item in claim_citations
        )
    for slide in presentation.slides:
        slide.notes_slide.notes_text_frame.text = source_notes
    presentation.core_properties.title = deck_title
    presentation.core_properties.subject = "Generated by the Multi-Agent Document and PPT POC"
    presentation.core_properties.author = "Supervisor agent workflow"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path


def content_from_profile(profile: DocumentProfile | dict[str, Any], request: str) -> dict[str, Any]:
    data = profile.to_dict() if isinstance(profile, DocumentProfile) else profile
    lines = [line.strip() for line in data.get("text", "").splitlines() if line.strip()]
    sections = []
    for item in data.get("sections", [])[:8]:
        heading = item.get("heading", "Context") if isinstance(item, dict) else str(item)
        sections.append({"heading": heading, "body": "Context carried forward from the source artifact.", "bullets": lines[:3]})
    if not sections:
        sections = [{"heading": "Source Context", "body": "The source artifact was converted into an editable deliverable.", "bullets": lines[:5]}]
    slides = [{"title": section["heading"], "bullets": section["bullets"] or [section["body"]]} for section in sections]
    while len(slides) < 11:
        slides.append({"title": f"Context extension {len(slides) + 1}", "bullets": ["Converted content remains editable and traceable."]})
    return {
        "title": f"Converted artifact: {data.get('file_name', 'source')}",
        "subtitle": request,
        "executive_summary": data.get("text", "")[:900] or "The source artifact did not contain extractable text.",
        "sections": sections,
        "slides": slides[:11],
        "sources": [],
    }
