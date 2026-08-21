from __future__ import annotations

from pathlib import Path
from typing import Any

from .parsers import parse_file


def _template_fidelity_checks(path: Path, file_type: str, profile: dict[str, Any]) -> list[dict[str, Any]]:
    source_file_type = profile.get("file_type") if isinstance(profile, dict) else None
    if source_file_type and source_file_type not in {file_type, "ppt" if file_type == "pptx" else file_type}:
        return []
    source_style = profile.get("style", {}) if isinstance(profile, dict) else {}
    if not source_style:
        return []
    generated = parse_file(path, path.name, ocr_enabled=False)
    generated_style = generated.style
    checks: list[dict[str, Any]] = []
    if file_type == "docx":
        source_margins = source_style.get("margins")
        generated_margins = generated_style.get("margins")
        if source_margins and generated_margins:
            checks.append({"name": "template_docx_margins", "passed": source_margins == generated_margins, "details": f"source={source_margins}, generated={generated_margins}"})
        source_fonts = set(source_style.get("font_names", []))
        generated_fonts = set(generated_style.get("font_names", []))
        checks.append({"name": "template_docx_fonts", "passed": bool(source_fonts.intersection(generated_fonts)), "details": f"shared={sorted(source_fonts.intersection(generated_fonts))}"})
        source_headers = source_style.get("headers_footers", {})
        generated_headers = generated_style.get("headers_footers", {})
        if source_headers:
            checks.append({"name": "template_docx_headers_footers", "passed": source_headers == generated_headers, "details": "Header/footer metadata preserved"})
        source_tables = source_style.get("tables", [])
        generated_tables = generated_style.get("tables", [])
        if source_tables:
            checks.append({"name": "template_docx_tables", "passed": len(generated_tables) >= len(source_tables), "details": f"source={len(source_tables)}, generated={len(generated_tables)}"})
    elif file_type == "pptx":
        source_layouts = set(source_style.get("layouts", []))
        generated_layouts = set(generated_style.get("layouts", []))
        source_masters = set(source_style.get("masters", []))
        generated_masters = set(generated_style.get("masters", []))
        if source_style.get("slide_size"):
            checks.append({"name": "template_pptx_slide_size", "passed": source_style.get("slide_size") == generated_style.get("slide_size"), "details": f"source={source_style.get('slide_size')}, generated={generated_style.get('slide_size')}"})
        if source_layouts:
            checks.append({"name": "template_pptx_layouts", "passed": source_layouts.issubset(generated_layouts), "details": f"preserved={sorted(source_layouts.intersection(generated_layouts))}"})
        if source_masters:
            checks.append({"name": "template_pptx_masters", "passed": source_masters.issubset(generated_masters), "details": f"preserved={sorted(source_masters.intersection(generated_masters))}"})
        source_shapes = {(item.get("name"), item.get("left"), item.get("top"), item.get("width"), item.get("height")) for item in source_style.get("shapes", [])}
        generated_shapes = {(item.get("name"), item.get("left"), item.get("top"), item.get("width"), item.get("height")) for item in generated_style.get("shapes", [])}
        if source_shapes:
            preserved = len(source_shapes.intersection(generated_shapes)) / len(source_shapes)
            checks.append({"name": "template_pptx_geometry", "passed": preserved >= 0.8, "details": f"preserved_ratio={preserved:.2f}"})
    return checks


def validate_artifact(path: Path, file_type: str, expected_slides: int | None = None, profile: dict[str, Any] | None = None) -> dict[str, Any]:
    checks: list[dict[str, Any]] = []
    try:
        if file_type == "docx":
            from docx import Document

            document = Document(path)
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
            checks.append({"name": "native_docx_readable", "passed": True, "details": f"{len(document.paragraphs)} paragraphs"})
            checks.append({"name": "editable_text_present", "passed": bool(text.strip()), "details": "Text paragraphs are present"})
        elif file_type == "pptx":
            from pptx import Presentation

            presentation = Presentation(path)
            text_shapes = sum(1 for slide in presentation.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
            checks.append({"name": "native_pptx_readable", "passed": True, "details": f"{len(presentation.slides)} slides"})
            checks.append({"name": "editable_text_present", "passed": text_shapes > 0, "details": f"{text_shapes} text shapes"})
            if expected_slides is not None:
                checks.append({"name": "expected_slide_count", "passed": len(presentation.slides) == expected_slides, "details": f"expected {expected_slides}, got {len(presentation.slides)}"})
        else:
            checks.append({"name": "file_exists", "passed": path.exists(), "details": str(path)})
        if profile:
            checks.extend(_template_fidelity_checks(path, file_type, profile))
    except Exception as exc:  # validation must return a report to the caller
        checks.append({"name": "artifact_open", "passed": False, "details": str(exc)})
    return {"passed": all(check["passed"] for check in checks), "checks": checks}
