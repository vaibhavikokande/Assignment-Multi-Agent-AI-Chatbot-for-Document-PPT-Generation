from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
import zipfile
from collections import Counter
from pathlib import Path
from typing import Any

from .models import DocumentProfile, new_id


SUPPORTED_EXTENSIONS = {".docx", ".pdf", ".ppt", ".pptx", ".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}


def _unique(values: list[Any], limit: int = 8) -> list[Any]:
    result: list[Any] = []
    for value in values:
        if value and value not in result:
            result.append(value)
    return result[:limit]


def _font_size(value: Any) -> float | None:
    try:
        return round(float(value.pt), 1) if value else None
    except (AttributeError, TypeError, ValueError):
        return None


def _rgb_color(value: Any) -> str | None:
    try:
        return str(value.rgb) if value and value.type is not None and value.rgb else None
    except (AttributeError, TypeError, ValueError):
        return None


def validate_file_signature(path: Path, file_name: str | None = None) -> None:
    """Reject extension/content mismatches before a file enters persistent storage."""
    path = Path(path)
    extension = Path(file_name or path.name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {extension or 'unknown'}")
    header = path.read_bytes()[:8]
    if extension in {".docx", ".pptx"} and not zipfile.is_zipfile(path):
        raise ValueError(f"Invalid {extension.lstrip('.').upper()} package")
    if extension == ".pdf" and not header.startswith(b"%PDF"):
        raise ValueError("Invalid PDF signature")
    if extension == ".ppt" and header != b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        raise ValueError("Invalid legacy PPT signature")
    if extension in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}:
        from PIL import Image

        try:
            with Image.open(path) as image:
                image.verify()
        except (OSError, ValueError) as exc:
            raise ValueError(f"Invalid image file: {exc}") from exc


def _ocr_file(path: Path) -> tuple[str, str]:
    """Use pytesseract when available, otherwise use the installed Tesseract CLI."""
    try:
        import pytesseract
        from PIL import Image

        text = pytesseract.image_to_string(Image.open(path)).strip()
        return text, "ok" if text else "empty_result"
    except ImportError:
        pass
    except (OSError, RuntimeError):
        return "", "unavailable"

    tesseract = os.getenv("TESSERACT_CMD") or shutil.which("tesseract")
    if not tesseract:
        return "", "unavailable"
    try:
        result = subprocess.run([tesseract, str(path), "stdout"], capture_output=True, text=True, timeout=60)
    except (OSError, subprocess.SubprocessError):
        return "", "unavailable"
    if result.returncode != 0:
        return "", "unavailable"
    text = result.stdout.strip()
    return text, "ok" if text else "empty_result"


def _parse_docx(path: Path, file_name: str) -> DocumentProfile:
    from docx import Document

    document = Document(path)
    paragraphs = [paragraph for paragraph in document.paragraphs if paragraph.text.strip()]
    text = "\n".join(paragraph.text.strip() for paragraph in paragraphs)
    font_names = []
    font_sizes = []
    for paragraph in paragraphs:
        for run in paragraph.runs:
            font_names.append(run.font.name)
            font_sizes.append(_font_size(run.font.size))
    sections = []
    for paragraph in paragraphs:
        style = paragraph.style.name if paragraph.style else "Normal"
        if "Heading" in style or paragraph.text.lower().startswith(("executive", "introduction", "overview", "conclusion")):
            sections.append({"heading": paragraph.text.strip(), "style": style})
    margins = []
    section_geometry = []
    for section in document.sections:
        geometry = {
            "top": round(section.top_margin.inches, 2),
            "bottom": round(section.bottom_margin.inches, 2),
            "left": round(section.left_margin.inches, 2),
            "right": round(section.right_margin.inches, 2),
            "width": round(section.page_width.inches, 2),
            "height": round(section.page_height.inches, 2),
            "orientation": str(section.orientation),
        }
        margins.append({key: geometry[key] for key in ("top", "bottom", "left", "right")})
        section_geometry.append(geometry)
    style_metadata = []
    for style in document.styles:
        if not getattr(style, "name", None):
            continue
        style_font = getattr(style, "font", None)
        style_metadata.append({
            "name": style.name,
            "type": str(getattr(style, "type", "")),
            "font_name": getattr(style_font, "name", None),
            "font_size": _font_size(getattr(style_font, "size", None)),
            "font_color": _rgb_color(getattr(style_font, "color", None)),
            "base_style": style.base_style.name if getattr(style, "base_style", None) else None,
        })
    paragraph_spacing = []
    for paragraph in paragraphs[:40]:
        paragraph_spacing.append({
            "style": paragraph.style.name if paragraph.style else "Normal",
            "before": round(paragraph.paragraph_format.space_before.pt, 1) if paragraph.paragraph_format.space_before else None,
            "after": round(paragraph.paragraph_format.space_after.pt, 1) if paragraph.paragraph_format.space_after else None,
            "line_spacing": str(paragraph.paragraph_format.line_spacing) if paragraph.paragraph_format.line_spacing else None,
            "alignment": str(paragraph.alignment) if paragraph.alignment else None,
        })
    header_footer = {
        "headers": ["\n".join(p.text for p in section.header.paragraphs if p.text.strip()) for section in document.sections],
        "footers": ["\n".join(p.text for p in section.footer.paragraphs if p.text.strip()) for section in document.sections],
    }
    table_metadata = []
    for table in document.tables[:20]:
        table_metadata.append({"rows": len(table.rows), "columns": len(table.columns), "style": table.style.name if table.style else None})
    return DocumentProfile(
        profile_id=new_id("profile"),
        file_name=file_name,
        file_type="docx",
        text=text,
        paragraph_count=len(paragraphs),
        table_count=len(document.tables),
        image_count=len(document.inline_shapes),
        style={
            "font_names": _unique(font_names) or ["Aptos"],
            "font_sizes": _unique(font_sizes) or [11.0],
            "heading_styles": _unique([section["style"] for section in sections]) or ["Heading 1"],
            "margins": margins[0] if margins else {"top": 0.7, "bottom": 0.7, "left": 0.8, "right": 0.8},
            "section_geometry": section_geometry,
            "styles": style_metadata[:80],
            "paragraph_spacing": paragraph_spacing,
            "headers_footers": header_footer,
            "tables": table_metadata,
            "custom_style_count": len(style_metadata),
        },
        sections=sections,
    )


def _parse_pdf(path: Path, file_name: str, ocr_enabled: bool = True) -> DocumentProfile:
    from pypdf import PdfReader

    reader = PdfReader(path)
    page_text = [(page.extract_text() or "").strip() for page in reader.pages]
    text = "\n\n".join(part for part in page_text if part)
    ocr_status = "not_required"
    if not text:
        ocr_status = "disabled" if not ocr_enabled else "unavailable"
        if ocr_enabled:
            pdftoppm = os.getenv("PDFTOPPM_PATH") or shutil.which("pdftoppm")
            if pdftoppm:
                temporary = Path(tempfile.mkdtemp(prefix="doc-pdf-ocr-"))
                prefix = temporary / "page"
                try:
                    result = subprocess.run([pdftoppm, "-png", "-r", "160", str(path), str(prefix)], capture_output=True, text=True, timeout=60)
                    if result.returncode == 0:
                        ocr_results = [_ocr_file(image) for image in sorted(temporary.glob("page-*.png"))]
                        ocr_parts = [part[0] for part in ocr_results if part[0]]
                        text = "\n\n".join(ocr_parts)
                        ocr_status = "ok" if text else (ocr_results[0][1] if ocr_results else "unavailable")
                except (OSError, subprocess.SubprocessError):
                    ocr_status = "unavailable"
                finally:
                    shutil.rmtree(temporary, ignore_errors=True)
    status = "ok" if text else "text_empty_scanned_content"
    return DocumentProfile(
        profile_id=new_id("profile"),
        file_name=file_name,
        file_type="pdf",
        text=text,
        page_count=len(reader.pages),
        paragraph_count=len([line for line in text.splitlines() if line.strip()]),
        style={"page_size": "A4", "page_text_lengths": [len(part) for part in page_text]},
        extraction_status=status,
        ocr_status=ocr_status,
    )


def _parse_pptx(path: Path, file_name: str) -> DocumentProfile:
    from pptx import Presentation

    presentation = Presentation(path)
    texts: list[str] = []
    font_names: list[str | None] = []
    font_sizes: list[float | None] = []
    layouts: list[str] = []
    shape_metadata: list[dict[str, Any]] = []
    backgrounds: list[dict[str, Any]] = []
    shape_count = 0
    for slide in presentation.slides:
        layouts.append(slide.slide_layout.name if slide.slide_layout else "Blank")
        background_color = None
        try:
            if slide.background.fill.type is not None:
                background_color = _rgb_color(slide.background.fill.fore_color)
        except (AttributeError, TypeError, ValueError):
            background_color = None
        backgrounds.append({"slide": slide.slide_id, "color": background_color})
        for shape in slide.shapes:
            shape_count += 1
            metadata = {
                "slide": slide.slide_id,
                "name": shape.name,
                "shape_type": int(shape.shape_type),
                "left": round(shape.left / 914400, 3),
                "top": round(shape.top / 914400, 3),
                "width": round(shape.width / 914400, 3),
                "height": round(shape.height / 914400, 3),
                "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
            }
            try:
                if shape.fill.type is not None:
                    metadata["fill_color"] = _rgb_color(shape.fill.fore_color)
            except (AttributeError, TypeError, ValueError):
                pass
            shape_metadata.append(metadata)
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    texts.append(paragraph.text.strip())
                for run in paragraph.runs:
                    font_names.append(run.font.name)
                    font_sizes.append(_font_size(run.font.size))
    text = "\n".join(texts)
    return DocumentProfile(
        profile_id=new_id("profile"),
        file_name=file_name,
        file_type="pptx",
        text=text,
        slide_count=len(presentation.slides),
        paragraph_count=len(texts),
        image_count=sum(1 for slide in presentation.slides for shape in slide.shapes if shape.shape_type == 13),
        style={
            "slide_size": {"width": round(presentation.slide_width / 914400, 2), "height": round(presentation.slide_height / 914400, 2)},
            "layouts": _unique(layouts),
            "font_names": _unique(font_names) or ["Aptos Display"],
            "font_sizes": _unique(font_sizes) or [24.0],
            "shape_count": shape_count,
            "masters": _unique([master.name for master in presentation.slide_masters]),
            "layout_metadata": [{"name": layout.name, "index": index, "placeholder_count": len(layout.placeholders)} for index, layout in enumerate(presentation.slide_layouts)],
            "backgrounds": backgrounds,
            "shapes": shape_metadata[:240],
            "theme": {"accent_colors": _unique([item.get("fill_color") for item in shape_metadata if item.get("fill_color")])},
        },
        sections=[{"heading": item, "slide": index + 1} for index, item in enumerate(texts[:12])],
    )


def _parse_image(path: Path, file_name: str, ocr_enabled: bool = True) -> DocumentProfile:
    from PIL import Image

    with Image.open(path) as image:
        width, height, mode = image.width, image.height, image.mode
    extracted = ""
    ocr_status = "disabled" if not ocr_enabled else "unavailable"
    if ocr_enabled:
        extracted, ocr_status = _ocr_file(path)
    return DocumentProfile(
        profile_id=new_id("profile"),
        file_name=file_name,
        file_type=path.suffix.lower().lstrip("."),
        text=extracted,
        image_count=1,
        style={"width": width, "height": height, "mode": mode},
        extraction_status="ok" if extracted else "image_metadata_only",
        ocr_status=ocr_status,
    )


def _convert_legacy_ppt(path: Path) -> Path:
    soffice = shutil.which("soffice") or os.getenv("SOFFICE_PATH")
    if not soffice:
        raise RuntimeError("Legacy .ppt input requires LibreOffice/soffice for conversion to .pptx")
    temporary = Path(tempfile.mkdtemp(prefix="doc-ppt-convert-"))
    result = subprocess.run([soffice, "--headless", "--convert-to", "pptx", "--outdir", str(temporary), str(path)], capture_output=True, text=True, timeout=60)
    converted = temporary / f"{path.stem}.pptx"
    if result.returncode != 0 or not converted.exists():
        raise RuntimeError(f"Legacy .ppt conversion failed: {result.stderr.strip() or result.stdout.strip()}")
    return converted


def parse_file(path: Path, file_name: str | None = None, ocr_enabled: bool = True) -> DocumentProfile:
    path = Path(path)
    display_name = file_name or path.name
    extension = path.suffix.lower()
    validate_file_signature(path, display_name)
    if extension == ".docx":
        return _parse_docx(path, display_name)
    if extension == ".pdf":
        return _parse_pdf(path, display_name, ocr_enabled)
    if extension == ".ppt":
        converted = _convert_legacy_ppt(path)
        try:
            profile = _parse_pptx(converted, display_name)
        finally:
            shutil.rmtree(converted.parent, ignore_errors=True)
        profile.file_type = "ppt"
        profile.extraction_status = "converted_to_pptx"
        return profile
    if extension == ".pptx":
        return _parse_pptx(path, display_name)
    return _parse_image(path, display_name, ocr_enabled)


def summarize_profile(profile: DocumentProfile) -> dict[str, Any]:
    return {
        "profile_id": profile.profile_id,
        "file_name": profile.file_name,
        "file_type": profile.file_type,
        "text_preview": profile.text[:500],
        "page_count": profile.page_count,
        "slide_count": profile.slide_count,
        "paragraph_count": profile.paragraph_count,
        "table_count": profile.table_count,
        "image_count": profile.image_count,
        "style": profile.style,
        "extraction_status": profile.extraction_status,
        "ocr_status": profile.ocr_status,
        "analysis": {
            "tone": "professional and structured" if profile.text else "visual template requiring content inference",
            "layout": "multi-page document" if profile.page_count else ("slide deck" if profile.slide_count else "image canvas"),
            "content_patterns": profile.sections[:8],
        },
    }
