from __future__ import annotations

from pathlib import Path


ROOT = Path(__file__).resolve().parent
INPUT_DIR = ROOT / "input"
OUTPUT_DIR = ROOT / "generated"


def create_sample_docx(path: Path | None = None) -> Path:
    from docx import Document
    from docx.shared import Inches, Pt

    path = path or INPUT_DIR / "Company_Proposal.docx"
    path.parent.mkdir(parents=True, exist_ok=True)
    document = Document()
    for section in document.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)
    normal = document.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    document.add_heading("Company Proposal", 0)
    document.add_paragraph("Practical, evidence-led transformation for enterprise teams.")
    document.add_heading("Overview", 1)
    document.add_paragraph("Our company helps organizations turn complex operating questions into measurable, governed programs.")
    document.add_heading("Approach", 1)
    document.add_paragraph("We combine discovery, workflow design, knowledge retrieval, evaluation, and change management.")
    document.add_paragraph("Evidence-led recommendations", style="List Bullet")
    document.add_paragraph("Reusable operating patterns", style="List Bullet")
    document.add_paragraph("Executive-ready communication", style="List Bullet")
    document.add_heading("Next Steps", 1)
    document.add_paragraph("Align on a priority workflow, evidence sources, owners, and success measures.")
    document.save(path)
    return path


def create_retail_ai_upload_template(path: Path | None = None) -> Path:
    """Create a polished DOCX template for the browser upload walkthrough."""
    from docx import Document
    from docx.enum.table import WD_ALIGN_VERTICAL
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor

    navy = RGBColor(11, 37, 69)
    blue = RGBColor(46, 116, 181)
    muted = RGBColor(89, 102, 120)
    pale_blue = "E8EEF5"
    light_gray = "F2F4F7"
    path = path or INPUT_DIR / "Northstar_Retail_AI_Enablement_Proposal_Template.docx"
    path.parent.mkdir(parents=True, exist_ok=True)
    document = Document()
    document.core_properties.title = "Northstar Retail AI Enablement Proposal Template"
    document.core_properties.subject = "Editable sample template for the Artifact Studio upload workflow"
    document.core_properties.author = "Artifact Studio"

    def set_run_font(run, *, size: float, color: RGBColor, bold: bool = False, italic: bool = False) -> None:
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:ascii"), "Calibri")
        run._element.rPr.rFonts.set(qn("w:hAnsi"), "Calibri")
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.bold = bold
        run.italic = italic

    def shade(cell, fill: str) -> None:
        properties = cell._tc.get_or_add_tcPr()
        shading = properties.find(qn("w:shd"))
        if shading is None:
            shading = OxmlElement("w:shd")
            properties.append(shading)
        shading.set(qn("w:fill"), fill)

    def set_cell_width(cell, width_dxa: int) -> None:
        properties = cell._tc.get_or_add_tcPr()
        width = properties.find(qn("w:tcW"))
        if width is None:
            width = OxmlElement("w:tcW")
            properties.append(width)
        width.set(qn("w:w"), str(width_dxa))
        width.set(qn("w:type"), "dxa")
        cell.width = Inches(width_dxa / 1440)

    def set_cell_margins(cell) -> None:
        properties = cell._tc.get_or_add_tcPr()
        margins = properties.find(qn("w:tcMar"))
        if margins is None:
            margins = OxmlElement("w:tcMar")
            properties.append(margins)
        for side, value in (("top", "80"), ("bottom", "80"), ("start", "120"), ("end", "120")):
            edge = margins.find(qn(f"w:{side}"))
            if edge is None:
                edge = OxmlElement(f"w:{side}")
                margins.append(edge)
            edge.set(qn("w:w"), value)
            edge.set(qn("w:type"), "dxa")

    def set_table_geometry(table, widths_dxa: list[int]) -> None:
        table.autofit = False
        properties = table._tbl.tblPr
        table_width = properties.find(qn("w:tblW"))
        if table_width is None:
            table_width = OxmlElement("w:tblW")
            properties.append(table_width)
        table_width.set(qn("w:w"), "9360")
        table_width.set(qn("w:type"), "dxa")
        indent = properties.find(qn("w:tblInd"))
        if indent is None:
            indent = OxmlElement("w:tblInd")
            properties.append(indent)
        indent.set(qn("w:w"), "120")
        indent.set(qn("w:type"), "dxa")
        layout = properties.find(qn("w:tblLayout"))
        if layout is None:
            layout = OxmlElement("w:tblLayout")
            properties.append(layout)
        layout.set(qn("w:type"), "fixed")
        for grid_column, width_dxa in zip(table._tbl.tblGrid.gridCol_lst, widths_dxa):
            grid_column.set(qn("w:w"), str(width_dxa))
        for row in table.rows:
            for cell, width_dxa in zip(row.cells, widths_dxa):
                set_cell_width(cell, width_dxa)
                set_cell_margins(cell)
                cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

    def add_rule() -> None:
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.space_after = Pt(12)
        border = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "12")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "2E74B5")
        border.append(bottom)
        paragraph._p.get_or_add_pPr().append(border)

    for section in document.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        section.header_distance = Inches(0.492)
        section.footer_distance = Inches(0.492)
        header = section.header.paragraphs[0]
        header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        header_run = header.add_run("NORTHSTAR RETAIL  |  STRATEGY TEMPLATE")
        set_run_font(header_run, size=8.5, color=muted, bold=True)
        footer = section.footer.paragraphs[0]
        footer.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        footer_run = footer.add_run("Northstar Retail  |  Retail AI Enablement")
        set_run_font(footer_run, size=8.5, color=muted)

    normal = document.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.font.color.rgb = navy
    normal.paragraph_format.space_after = Pt(8)
    normal.paragraph_format.line_spacing = 1.333
    for style_name, size, color, before, after in (("Heading 1", 16, blue, 18, 10), ("Heading 2", 13, blue, 12, 6), ("Heading 3", 12, navy, 8, 4)):
        style = document.styles[style_name]
        style.font.name = "Calibri"
        style.font.size = Pt(size)
        style.font.color.rgb = color
        style.font.bold = True
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(after)

    kicker = document.add_paragraph()
    kicker.alignment = WD_ALIGN_PARAGRAPH.CENTER
    kicker.paragraph_format.space_before = Pt(36)
    kicker.paragraph_format.space_after = Pt(8)
    set_run_font(kicker.add_run("NORTHSTAR RETAIL"), size=12, color=muted, bold=True)
    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.paragraph_format.space_after = Pt(6)
    set_run_font(title.add_run("Retail AI Enablement Proposal"), size=26, color=navy, bold=True)
    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.paragraph_format.space_after = Pt(24)
    set_run_font(subtitle.add_run("A practical roadmap for customer service, store operations, and inventory decisions"), size=13, color=muted)
    add_rule()

    metadata = document.add_table(rows=2, cols=2)
    metadata.style = "Table Grid"
    set_table_geometry(metadata, [4680, 4680])
    metadata_rows = (("Prepared for", "Northstar Retail Leadership Team"), ("Planning horizon", "90-day pilot and scale recommendation"), ("Prepared by", "Strategy & Transformation Office"), ("Decision requested", "Approve discovery and pilot design"))
    for cell, (label, value) in zip((cell for row in metadata.rows for cell in row.cells), metadata_rows):
        shade(cell, light_gray)
        paragraph = cell.paragraphs[0]
        paragraph.paragraph_format.space_after = Pt(1)
        set_run_font(paragraph.add_run(f"{label}: "), size=9.5, color=navy, bold=True)
        set_run_font(paragraph.add_run(value), size=9.5, color=navy)

    document.add_heading("Executive context", 1)
    document.add_paragraph("Northstar Retail is evaluating targeted AI capabilities to improve customer responsiveness, simplify store-team work, and create more confident inventory decisions. This template is intentionally concise, evidence-led, and oriented toward measurable outcomes.")
    document.add_heading("Priority outcomes", 2)
    for bullet in ("Reduce time to resolve routine customer-service questions.", "Give store teams clear, policy-aligned operational guidance.", "Improve planner visibility into demand signals and inventory exceptions."):
        document.add_paragraph(bullet, style="List Bullet")

    document.add_heading("90-day roadmap", 1)
    roadmap = document.add_table(rows=1, cols=3)
    roadmap.style = "Table Grid"
    set_table_geometry(roadmap, [2160, 4320, 2880])
    for cell, value in zip(roadmap.rows[0].cells, ("Phase", "Workstream", "Success signal")):
        shade(cell, pale_blue)
        paragraph = cell.paragraphs[0]
        paragraph.paragraph_format.space_after = Pt(0)
        set_run_font(paragraph.add_run(value), size=10, color=navy, bold=True)
    for row_values in (("Weeks 1–3", "Map high-volume service, store, and planning decisions.", "Prioritized use-case backlog and baseline metrics."), ("Weeks 4–7", "Prototype grounded assistant workflows with human review.", "Pilot users complete tasks faster with documented controls."), ("Weeks 8–12", "Measure value, adoption, risk, and scale requirements.", "Leadership decision on rollout investment and operating model.")):
        cells = roadmap.add_row().cells
        for cell, value in zip(cells, row_values):
            paragraph = cell.paragraphs[0]
            paragraph.paragraph_format.space_after = Pt(0)
            set_run_font(paragraph.add_run(value), size=9.5, color=navy)

    document.add_page_break()
    document.add_heading("Recommended pilot scope", 1)
    document.add_paragraph("Start with a small number of repeatable workflows where the business owner, source information, and success measures are already clear. The pilot should combine retrieval-grounded guidance with explicit human review for high-impact decisions.")
    document.add_heading("Design principles", 2)
    for bullet in ("Use approved internal knowledge alongside current external evidence.", "Maintain citations, run records, and version history for every deliverable.", "Design the experience for store and service teams, not only technical users.", "Measure quality, adoption, time saved, and risk before scale-up."):
        document.add_paragraph(bullet, style="List Bullet")
    document.add_heading("Decision checklist", 1)
    for step in ("Confirm one executive sponsor and accountable workflow owner.", "Approve the pilot use cases, source systems, and review controls.", "Set baseline metrics and a 90-day decision date."):
        document.add_paragraph(step, style="List Number")
    document.add_heading("Template guidance", 1)
    document.add_paragraph("When using this file in Artifact Studio, ask the supervisor to retain the concise executive tone, navy-and-blue hierarchy, structured roadmap, and measurable-outcome framing. The generated DOCX and PPTX remain editable for review and iteration.")
    document.save(path)
    return path


def create_sample_pptx(path: Path | None = None) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    path = path or INPUT_DIR / "Company_Template.pptx"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    for slide in list(presentation.slides):
        slide_id = slide.slide_id
        for item in list(presentation.slides._sldIdLst):
            if item.id == slide_id:
                presentation.part.drop_rel(item.rId)
                presentation.slides._sldIdLst.remove(item)
    for index, (title_text, body_text) in enumerate([
        ("Company Template", "Executive narrative and decision support"),
        ("A practical point of view", "Evidence-led\nGoverned\nMeasurable"),
        ("The next move", "Start with one priority workflow and scale what works."),
    ]):
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
        body = slide.placeholders[1]
        body.text = body_text
        for paragraph in body.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
    presentation.core_properties.title = "Company Presentation Template"
    presentation.save(path)
    return path


def create_samples() -> dict[str, str]:
    doc = create_sample_docx()
    ppt = create_sample_pptx()
    context = INPUT_DIR / "enterprise_context.md"
    context.write_text("# Enterprise context\n\nThe company prioritizes secure adoption, measurable value, evidence-led recommendations, and reusable workflow patterns.\n", encoding="utf-8")
    return {"docx": str(doc), "pptx": str(ppt), "context": str(context)}


if __name__ == "__main__":
    for key, value in create_samples().items():
        print(f"{key}: {value}")
